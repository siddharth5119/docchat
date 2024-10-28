import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
#spliting the pdf vector to chunks
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
# embeddings for google
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
#help in chat
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import csv
import docx  # for handling DOCX files
import pptx  # for handling PPTX files
from bs4 import BeautifulSoup  # for handling HTML files


load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))


# def get_pdf_text(pdf_docs):
#     text=""
#     for pdf in pdf_docs:
#         pdf_reader= PdfReader(pdf)
#         for page in pdf_reader.pages:
#             text+= page.extract_text()
#     return  text


def get_document_text(documents):
    text = ""
    for doc in documents:
        file_type = os.path.splitext(doc.name)[1].lower()

        if file_type == '.pdf':
            pdf_reader = PdfReader(doc)
            for page in pdf_reader.pages:
                text += page.extract_text()

        elif file_type == '.txt':
            text += doc.read().decode("utf-8")

        elif file_type == '.csv':
            csv_reader = csv.reader(doc.read().decode("utf-8").splitlines())
            for row in csv_reader:
                text += ' '.join(row) + '\n'

        elif file_type == '.docx':
            docx_reader = docx.Document(doc)
            for para in docx_reader.paragraphs:
                text += para.text + '\n'

        elif file_type == '.pptx':
            ppt = pptx.Presentation(doc)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'

        elif file_type == '.html':
            soup = BeautifulSoup(doc, "html.parser")
            text += soup.get_text()

        else:
            st.warning(f"Unsupported file type: {file_type}")
    
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")


def get_conversational_chain():

    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro",
                             temperature=0.3)

    prompt = PromptTemplate(template = prompt_template, input_variables = ["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)

    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")
    
    new_db = FAISS.load_local("faiss_index", embeddings)
    docs = new_db.similarity_search(user_question)

    chain = get_conversational_chain()

    
    response = chain(
        {"input_documents":docs, "question": user_question}
        , return_only_outputs=True)

    print(response)
    st.write("Reply: ", response["output_text"])


# # Use the function within your app
# uploaded_files = st.file_uploader("Upload files", type=["pdf", "txt", "csv", "docx", "pptx", "html"], accept_multiple_files=True)
# if uploaded_files:
#     document_text = get_document_text(uploaded_files)
#     st.write(document_text)


