import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import csv
import docx  # for handling DOCX files
import pptx  # for handling PPTX files
from bs4 import BeautifulSoup  # for handling HTML files

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Function to extract text from various document types
def get_document_text(documents):
    text = ""
    for doc in documents:
        file_type = os.path.splitext(doc.name)[1].lower()

        if file_type == '.pdf':
            pdf_reader = PdfReader(doc)
            for page in pdf_reader.pages:
                text += page.extract_text()

        elif file_type == '.txt':
            text += doc.read().decode("utf-8")

        elif file_type == '.csv':
            csv_reader = csv.reader(doc.read().decode("utf-8").splitlines())
            for row in csv_reader:
                text += ' '.join(row) + '\n'

        elif file_type == '.docx':
            docx_reader = docx.Document(doc)
            for para in docx_reader.paragraphs:
                text += para.text + '\n'

        elif file_type == '.pptx':
            ppt = pptx.Presentation(doc)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'

        elif file_type == '.html':
            soup = BeautifulSoup(doc, "html.parser")
            text += soup.get_text()

        else:
            st.warning(f"Unsupported file type: {file_type}")
    
    return text

# Split text into chunks for vector storage
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

# Create a vector store using FAISS
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

# Load conversational chain for Q&A
def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. Make sure to provide all the details.
    If the answer is not in the provided context, just say, "answer is not available in the context."
    Don't provide wrong answers.

    Context:\n {context}?\n
    Question:\n {question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Handle user input, query the FAISS index, and return a response
def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    
    # Allow dangerous deserialization if you trust the file source
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()

    response = chain(
        {"input_documents": docs, "question": user_question}, return_only_outputs=True
    )

    st.write("Reply: ", response["output_text"])

# Streamlit UI
st.set_page_config(page_title="Document Q&A System", layout="wide")

# Header section with styling
st.markdown("<h1 style='text-align: center; color: #3498db;'>📝 Document Q&A System</h1>", unsafe_allow_html=True)
st.markdown("<p style='text-align: center; color: #2ecc71;'>Upload your documents, preview the content, and ask questions!</p>", unsafe_allow_html=True)

# File uploader for different file types
uploaded_files = st.file_uploader("Upload files", type=["pdf", "txt", "csv", "docx", "pptx", "html"], accept_multiple_files=True)

# If files are uploaded, process them
if uploaded_files:
    st.markdown("<h3 style='color: #9b59b6;'>📄 Extracted Content Previews:</h3>", unsafe_allow_html=True)

    # Document content previews
    for doc in uploaded_files:
        with st.expander(f"Preview: {doc.name}"):
            document_text = get_document_text([doc])
            
            # Set character limit for preview
            char_limit = 1000  # Limit to 1000 characters
            if len(document_text) > char_limit:
                preview_text = document_text[:char_limit] + " ... (Read More)"
                st.write(preview_text)
                
                # Button to show full text
                if st.button(f"Show Full Text for {doc.name}"):
                    st.write(document_text)  # Show the entire document text
            else:
                st.write(document_text)  # Show the entire text if it's under the limit

    # Extract text and generate text chunks for all uploaded documents
    all_text = get_document_text(uploaded_files)
    text_chunks = get_text_chunks(all_text)
    get_vector_store(text_chunks)  # Store in FAISS

    # Ask a question section
    st.markdown("<h3 style='color: #e74c3c;'>❓ Ask a Question Based on the Document:</h3>", unsafe_allow_html=True)

    user_question = st.text_input("Enter your question here")
    if user_question:
        user_input(user_question)
