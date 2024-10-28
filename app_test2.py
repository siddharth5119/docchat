import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import csv
import docx  # for handling DOCX files
import pptx  # for handling PPTX files
from bs4 import BeautifulSoup  # for handling HTML files

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Initialize session state for conversation
if 'conversation' not in st.session_state:
    st.session_state['conversation'] = []
if 'topics' not in st.session_state:
    st.session_state['topics'] = {}

# Function to extract text from various document types
def get_document_text(doc):
    text = ""
    file_type = os.path.splitext(doc.name)[1].lower()

    if file_type == '.pdf':
        pdf_reader = PdfReader(doc)
        for page in pdf_reader.pages:
            text += page.extract_text()

    elif file_type == '.txt':
        text += doc.read().decode("utf-8")

    elif file_type == '.csv':
        csv_reader = csv.reader(doc.read().decode("utf-8").splitlines())
        for row in csv_reader:
            text += ' '.join(row) + '\n'

    elif file_type == '.docx':
        docx_reader = docx.Document(doc)
        for para in docx_reader.paragraphs:
            text += para.text + '\n'

    elif file_type == '.pptx':
        ppt = pptx.Presentation(doc)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'

    elif file_type == '.html':
        soup = BeautifulSoup(doc, "html.parser")
        text += soup.get_text()

    else:
        st.warning(f"Unsupported file type: {file_type}")

    return text

# Split text into chunks for vector storage
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

# Create a vector store using FAISS
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

# Load conversational chain for Q&A
def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. Make sure to provide all the details.
    If the answer is not in the provided context, just say, "Answer is not available in the context."
    Don't provide wrong answers.

    Context:\n {context}?\n
    Question:\n {question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Handle user input, query the FAISS index, and return a response with references
def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    
    # Search for similar documents and run the Q&A chain
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain(
        {"input_documents": docs, "question": user_question}, return_only_outputs=True
    )
    
    # Collect references for each document
    references = []
    for doc in docs:
        references.append(f"Document: {doc.metadata.get('document_name', 'Unknown')}, Page: {doc.metadata.get('page_number', 'N/A')}")

    # Store the question and answer in the conversation history
    st.session_state['conversation'].append({"question": user_question, "answer": response["output_text"], "references": references})

# Streamlit UI for conversation history
def display_conversation():
    for idx, convo in enumerate(st.session_state['conversation']):
        st.write(f"**Q{idx+1}: {convo['question']}**")
        st.write(f"**A{idx+1}: {convo['answer']}**")
        st.write(f"_References:_ {', '.join(convo['references'])}")
        st.markdown("---")

# Streamlit Sidebar for file uploads
st.sidebar.title("📂 Upload Documents")
uploaded_files = st.sidebar.file_uploader("Upload files", type=["pdf", "txt", "csv", "docx", "pptx", "html"], accept_multiple_files=True)

if uploaded_files:
    for file in uploaded_files:
        # Display the preview button for each file
        if st.sidebar.button(f"Preview {file.name}"):
            document_text = get_document_text(file)
            with st.expander(f"📄 Preview: {file.name} (Full Text)"):
                st.write(document_text)  # Display full document preview

    st.sidebar.markdown("---")

    # If files are uploaded, process and split text
    combined_text = ""
    for file in uploaded_files:
        combined_text += get_document_text(file)

    text_chunks = get_text_chunks(combined_text)
    get_vector_store(text_chunks)  # Store in FAISS

# Main chat interface
st.title("🗨️ Document Q&A Chat")
st.markdown("Ask questions about the uploaded documents and get detailed responses with references!")

# Input box for questions
user_question = st.text_input("Type your question here")

if st.button("Submit Question") and user_question:
    user_input(user_question)

# Display the conversation history
st.markdown("### 💬 Conversation History")
display_conversation()
